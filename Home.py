from __future__ import annotations
import os, json, base64, html, re, requests, mimetypes
from pathlib import Path
from typing import Dict, List, Optional, Tuple
import streamlit as st

# =======================
# ❖ App / Assets
# =======================
APP_TITLE = "TomTom Tax Agent"
APP_ICON = ".streamlit/TomTom-Logo.png"
LOGO_PATH = Path(".streamlit/TomTom-Logo.png")

# Endpoint config
LLM_API_KEY = os.getenv("AZURE_API_KEY", "")
LLM_ENDPOINT = os.getenv("AZURE_API_ENDPOINT", "")
if LLM_ENDPOINT and not LLM_ENDPOINT.lower().startswith("http"):
    raise RuntimeError("AZURE_API_ENDPOINT must be a full https URL to your Prompt Flow endpoint.")
if LLM_API_KEY and LLM_API_KEY.lower().startswith("http"):
    raise RuntimeError("AZURE_API_KEY must be the secret key, not a URL.")

# Session keys
SK_MSGS = "messages"
ATTACH_CTX_KEY = "attachment_contexts"
KB_STATE_KEY = "BASE_KB_CHUNKS"   # preloaded local KB chunks

MAX_CONTEXT_MESSAGES = 20
PF_MAX_TURNS = 10

# =======================
# ❖ Bootstrap KB (local)
# =======================
BOOTSTRAP_DOCS_DIR = "kb"          # put your PDFs/DOCX/TXT/etc here
MAX_BOOTSTRAP_TOTAL_CHARS = 200_000
KB_CHUNK_SIZE = 1200
KB_CHUNK_OVERLAP = 150
KB_TOP_K = 6

# =======================
# ❖ Optional libs
# =======================
try:
    from PIL import Image
    PIL_AVAILABLE = True
except Exception:
    PIL_AVAILABLE = False

# ============ Page setup (keep light theme from your config) ============
if PIL_AVAILABLE and Path(APP_ICON).exists():
    try:
        _icon_obj = Image.open(APP_ICON)
    except Exception:
        _icon_obj = APP_ICON
else:
    _icon_obj = APP_ICON

st.set_page_config(
    page_title=APP_TITLE,
    page_icon=_icon_obj,
    initial_sidebar_state="expanded",
)

# Minimal CSS (keep light theme; TomTom red accents)
st.markdown(
    """
    <style>
      [data-testid="stSidebarNav"], section[data-testid="stSidebar"] nav { display: none !important; }

      /* Hide built-in chat avatars (we render our own emoji/icons) */
      [data-testid="chatAvatarIcon-user"], [data-testid="chatAvatarIcon-assistant"] { display: none !important; }

      /* Chat margins */
      .chat-margin-container { margin-left: 100px; margin-right: 100px; }
      @media (max-width: 900px) {
        .chat-margin-container { margin-left: 10px; margin-right: 10px; }
      }

      /* TomTom red accents */
      .stButton>button {
        background-color: #D30000 !important;
        color: #fff !important;
        border: none !important;
        border-radius: 9999px !important;
      }
      .stButton>button:hover { filter: brightness(1.1); }
      h3, h4 { color: #D30000 !important; }
      hr { border: none; border-top: 2px solid #D30000 !important; margin: .5rem 0; }
      ::-webkit-scrollbar-thumb { background: #D30000 !important; border-radius: 4px; }

      /* Simple bubbles, no shadow */
      .msg-row { display:flex; align-items:flex-start; width:100%; margin-bottom: 0.85em; }
      .msg-left { justify-content:flex-start; }
      .msg-right { justify-content:flex-end; }
      .msg-bubble { max-width:75%; padding:.6em .9em; border-radius:14px; line-height:1.25; font-size:1.02em; }
      .msg-user  { background: rgba(0,0,0,0.06); }
      .msg-assistant { background: rgba(0,0,0,0.04); }
      .msg-icon { font-size: 28px; margin: 0 10px; line-height:1; }

      /* Floating attach popover trigger at bottom-left area (adjust if needed) */
      .attach-anchor { position: fixed; left: 100px; bottom: 68px; z-index: 9999; }
      @media (max-width: 900px) { .attach-anchor { left: 10px; bottom: 68px; } }
    </style>
    """,
    unsafe_allow_html=True,
)

# =======================
# ❖ Logo helper
# =======================
def logo_img_base64() -> Optional[str]:
    if LOGO_PATH and LOGO_PATH.exists():
        try:
            return base64.b64encode(LOGO_PATH.read_bytes()).decode("utf-8")
        except Exception:
            return None
    return None

def show_logo(center: bool = True, max_width: str = "180px") -> None:
    b64 = logo_img_base64()
    if b64:
        align = "margin-left:auto;margin-right:auto;" if center else ""
        st.markdown(
            f"""
            <div style="text-align:center; margin: 0 0 .25rem 0;">
                <img src="data:image/png;base64,{b64}" style="max-width:{max_width}; {align}" />
            </div>
            """,
            unsafe_allow_html=True,
        )

# =======================
# ❖ Countries (sidebar)
# =======================
COUNTRIES: List[Tuple[str, str]] = [
    ("Netherlands", "NL"), ("Belgium", "BE"), ("Germany", "DE"), ("France", "FR"),
    ("United Kingdom", "GB"), ("Spain", "ES"), ("Italy", "IT"), ("Poland", "PL"),
    ("United States", "US"), ("Canada", "CA"),
]
def find_country_by_name(name: str) -> Optional[Tuple[str, str]]:
    for n, c in COUNTRIES:
        if n == name:
            return (n, c)
    return None

# =======================
# ❖ PF chat history (question/answer pairs)
# =======================
def to_pf_chat_history(streamlit_msgs: List[Dict[str, str]], max_pairs: int = PF_MAX_TURNS) -> List[Dict]:
    """
    Convert Streamlit chat list -> Prompt Flow turns:
      [{"inputs":{"question": <user>}, "outputs":{"answer": <assistant>}}, ...]
    """
    pairs: List[Dict] = []
    pending_user: Optional[str] = None
    if not isinstance(streamlit_msgs, list) or not streamlit_msgs:
        return pairs

    for m in streamlit_msgs:
        role = (m.get("role") or "").strip().lower()
        content = (m.get("content") or "").strip()
        if not content:
            continue
        if role == "user":
            pending_user = content
        elif role == "assistant" and pending_user is not None:
            pairs.append({"inputs": {"question": pending_user}, "outputs": {"answer": content}})
            pending_user = None

    return pairs[-max_pairs:]

# =======================
# ❖ Default-output filter (AML tutorial guard)
# =======================
DEFAULT_SIGNATURES = [
    "ml_client.compute.begin_create_or_update(compute_instance)",
    "from azure.ai.ml import mlclient",
    "from azure.ai.ml.entities import computeinstance",
    "from azure.identity import defaultazurecredential",
    "standard_ds3_v2",
    "steps to create a compute instance",
    "to create an azure machine learning compute instance",
    "pip install azure-ai-ml",
    "computeinstance(",
    "begin_create_or_update(",
    "compute instance '",
    "azure machine learning workspace",
    "vm size (e.g., standard_ds3_v2)",
    "monitor the creation process",
]
def looks_like_default(text: str) -> bool:
    t = (text or "").strip().lower()
    if not t:
        return False
    if any(sig in t for sig in DEFAULT_SIGNATURES):
        return True
    if ("compute instance" in t and ("azureml" in t or "azure.ai.ml" in t)):
        return True
    if "###" in t and "```" in t and "pip install azure-ai-ml" in t:
        return True
    return False

# =======================
# ❖ Upload extraction (many types)
# =======================
def _bytes_to_text(b: bytes) -> str:
    try:
        return b.decode("utf-8", errors="ignore")
    except Exception:
        return ""

def _pdf_to_text(b: bytes) -> str:
    try:
        import pypdf
        from io import BytesIO
        r = pypdf.PdfReader(BytesIO(b))
        return "\n".join((p.extract_text() or "") for p in r.pages)
    except Exception:
        return ""

def _docx_to_text(b: bytes) -> str:
    try:
        from io import BytesIO
        from docx import Document
        doc = Document(BytesIO(b))
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except Exception:
        return ""

def _xlsx_to_text(b: bytes) -> str:
    try:
        import pandas as pd
        from io import BytesIO
        xls = pd.ExcelFile(BytesIO(b))
        parts = []
        for sheet in xls.sheet_names:
            try:
                df = xls.parse(sheet, dtype=str).fillna("")
                parts.append(f"== Sheet: {sheet} ==\n" + df.to_csv(index=False))
            except Exception:
                continue
        return "\n\n".join(parts)
    except Exception:
        return ""

def _pptx_to_text(b: bytes) -> str:
    try:
        from io import BytesIO
        from pptx import Presentation
        prs = Presentation(BytesIO(b))
        texts = []
        for i, slide in enumerate(prs.slides, 1):
            buf = [f"== Slide {i} =="]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    buf.append(shape.text)
            texts.append("\n".join(buf))
        return "\n\n".join(texts)
    except Exception:
        return ""

def _rtf_to_text(b: bytes) -> str:
    t = _bytes_to_text(b)
    t = re.sub(r"\\[a-zA-Z]+\d*", "", t)
    t = re.sub(r"[{}]", "", t)
    return t

def _html_to_text(b: bytes) -> str:
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(_bytes_to_text(b), "html.parser")
        return soup.get_text("\n")
    except Exception:
        return _bytes_to_text(b)

def _xml_to_text(b: bytes) -> str:
    return _bytes_to_text(b)

def _image_to_text(b: bytes) -> str:
    # Optional OCR if pytesseract + Tesseract are installed
    try:
        from PIL import Image
        import pytesseract
        from io import BytesIO
        img = Image.open(BytesIO(b))
        return pytesseract.image_to_string(img) or ""
    except Exception:
        return ""

MAX_TOTAL_ATTACHMENT_CHARS = 60_000
MAX_PER_FILE_CHARS = 20_000

def extract_text_from_upload(uploaded_file) -> str:
    name = (uploaded_file.name or "").lower()
    try:
        data = uploaded_file.getvalue() if hasattr(uploaded_file, "getvalue") else uploaded_file.read()
    except Exception:
        data = b""
    if not data:
        return ""

    text = ""
    if name.endswith(".pdf"):
        text = _pdf_to_text(data)
    elif name.endswith(".docx"):
        text = _docx_to_text(data)
    elif name.endswith((".txt", ".md", ".csv", ".tsv", ".json", ".yaml", ".yml", ".log")):
        text = _bytes_to_text(data)
    elif name.endswith(".xlsx"):
        text = _xlsx_to_text(data)
    elif name.endswith(".pptx"):
        text = _pptx_to_text(data)
    elif name.endswith((".html", ".htm")):
        text = _html_to_text(data)
    elif name.endswith(".rtf"):
        text = _rtf_to_text(data)
    elif name.endswith((".xml", )):
        text = _xml_to_text(data)
    elif name.endswith((".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tif", ".tiff")):
        text = _image_to_text(data) or f"[Image uploaded: {uploaded_file.name}]"
    else:
        guessed, _ = mimetypes.guess_type(name)
        if guessed and "text" in guessed:
            text = _bytes_to_text(data)
        else:
            text = _bytes_to_text(data)

    text = (text or "").strip()
    if len(text) > MAX_PER_FILE_CHARS:
        text = text[:MAX_PER_FILE_CHARS] + "\n...[truncated]"

    if text:
        return f"=== FILE: {uploaded_file.name} ===\n{str(text)}"
    return ""

# =======================
# ❖ Local KB reader / chunker
# =======================
def _chunk_text(text: str, size: int = KB_CHUNK_SIZE, overlap: int = KB_CHUNK_OVERLAP) -> List[str]:
    text = (text or "").strip()
    chunks = []
    if not text:
        return chunks
    start = 0
    n = len(text)
    while start < n:
        end = min(start + size, n)
        chunks.append(text[start:end])
        if end == n:
            break
        start = max(0, end - overlap)
    return chunks

def read_local_file_to_text(path: Path) -> str:
    name = path.name.lower()
    try:
        data = path.read_bytes()
    except Exception:
        return ""
    fake_uploaded = type("F", (), {"name": path.name, "read": lambda self=None: data, "getvalue": lambda self=None: data})
    return extract_text_from_upload(fake_uploaded)

def build_bootstrap_kb(dir_path: str) -> List[str]:
    folder = Path(dir_path)
    if not folder.exists() or not folder.is_dir():
        return []
    chunks: List[str] = []
    for p in folder.rglob("*"):
        if not p.is_file():
            continue
        try:
            txt = read_local_file_to_text(p)
            if not txt:
                continue
            prefixed = f"=== FILE: {p.name} ===\n{txt}"
            for ch in _chunk_text(prefixed):
                chunks.append(ch)
        except Exception:
            continue

    joined = "\n\n".join(chunks)
    if len(joined) > MAX_BOOTSTRAP_TOTAL_CHARS:
        joined = joined[:MAX_BOOTSTRAP_TOTAL_CHARS] + "\n...[KB truncated]"
        return joined.split("\n\n")
    return chunks

def _tokenize(s: str) -> List[str]:
    return re.findall(r"[a-z0-9]+", (s or "").lower())

def retrieve_kb_chunks(query: str, kb_chunks: List[str], top_k: int = KB_TOP_K) -> List[str]:
    q = set(_tokenize(query))
    if not q or not kb_chunks:
        return []
    scored = []
    for ch in kb_chunks:
        c = set(_tokenize(ch))
        score = len(q.intersection(c))
        if score > 0:
            scored.append((score, ch))
    scored.sort(key=lambda x: x[0], reverse=True)
    return [ch for _, ch in scored[:top_k]]

# =======================
# ❖ LLM call (Prompt Flow)
# =======================
def get_llm_response(prompt: str, *_ignore_context: str) -> str:
    if not LLM_API_KEY or not LLM_ENDPOINT:
        return "⚠️ Missing LLM configuration. Set AZURE_API_KEY and AZURE_API_ENDPOINT."

    endpoint = LLM_ENDPOINT.rstrip("/")
    urls = [endpoint] + ([endpoint + "/score"] if not endpoint.endswith("/score") else [])
    is_aml = ".inference.ml.azure.com" in endpoint.lower()

    def headers_bearer():
        return {"Content-Type": "application/json", "Authorization": f"Bearer {LLM_API_KEY}"}
    def headers_apikey():
        return {"Content-Type": "application/json", "api-key": LLM_API_KEY}

    headers_list = [headers_bearer(), headers_apikey()]  # try both

    # PF-style history (question/answer)
    history_pf = to_pf_chat_history(st.session_state.get(SK_MSGS, []))
    if not history_pf:
        history_pf = [{"inputs": {"question": "Hi"}, "outputs": {"answer": "Hello!"}}]

    # Build prompt with KB + attachments
    kb_chunks = st.session_state.get(KB_STATE_KEY, []) or []
    kb_hits = retrieve_kb_chunks(prompt, kb_chunks, top_k=KB_TOP_K)
    kb_text = "\n\n".join(kb_hits).strip()

    attach_blocks = st.session_state.get(ATTACH_CTX_KEY, []) or []
    attach_text = "\n\n".join(attach_blocks).strip()

    sections = []
    if kb_text:
        sections.append(f"KB_CONTEXT:\n{kb_text}")
    if attach_text:
        sections.append(f"ATTACHMENT_CONTEXT:\n{attach_text}")
    sections.append(f"USER_QUESTION:\n{prompt}")
    combined_prompt = "\n\n".join(sections)

    body_foundry = {"inputs": {"question": combined_prompt, "chat_history": history_pf}}
    body_aml     = {"input_data": {"inputs": {"question": combined_prompt, "chat_history": history_pf}}}
    payloads = [("foundry.inputs", body_foundry), ("aml.input_data.inputs", body_aml), ("raw.top", {"question": combined_prompt, "chat_history": history_pf})]

    last_status = last_text = last_where = None

    for url in urls:
        for headers in headers_list:
            for tag, body in payloads:
                where = f"{url} [{ 'Bearer' if 'Authorization' in headers else 'api-key' } | {tag}]"
                try:
                    resp = requests.post(url, headers=headers, json=body, timeout=120)
                    last_status, last_text, last_where = resp.status_code, resp.text, where
                    if resp.status_code != 200:
                        continue

                    # Try JSON first
                    try:
                        data = resp.json()
                    except Exception:
                        txt = (resp.text or "").strip()
                        if txt and looks_like_default(txt):
                            continue
                        if txt:
                            return txt
                        continue

                    # Preferred extraction paths for this flow
                    content = (
                        (data.get("outputs") or {}).get("answer")
                        or data.get("answer")
                        or data.get("output")
                        or data.get("result")
                        or data.get("prediction")
                        or data.get("value")
                    )
                    if content:
                        if looks_like_default(content):
                            continue
                        return content

                    # fallback stringify
                    cand = json.dumps(data, ensure_ascii=False)
                    if cand and looks_like_default(cand):
                        continue
                    if cand:
                        return cand[:4000]

                except requests.RequestException as e:
                    last_text = f"Network error: {e}"
                    continue

    debug = f"Last attempt: {last_where} → HTTP {last_status}\nBody (first 400 chars):\n{(last_text or '')[:400]}"
    return (
        "⚠️ The endpoint responded with a sample/default output or an error. "
        "This usually means the request schema or auth header doesn't match the deployed flow. "
        "Open the endpoint's **Test** tab, run a test, click **View request**, and align the JSON.\n\n" + debug
    )

# =======================
# ❖ Simple render helpers
# =======================
def format_llm_reply_to_html(md: str) -> str:
    """
    Lightweight Markdown -> HTML conversion for LLM outputs.
    Supports: headings, bold/italic, lists, tables, horizontal rules, paragraphs.
    """
    if not md:
        return ""

    md = md.strip()

    # --- Normalization ---
    md = re.sub(r"\r\n", "\n", md)   # unify line breaks
    md = re.sub(r"\n{3,}", "\n\n", md)  # collapse 3+ newlines → 2

    # --- Inline formatting ---
    def fmt_inline(s: str) -> str:
        s = html.escape(s)
        s = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", s)   # bold
        s = re.sub(r"(?<!\*)\*(?!\s)(.+?)(?<!\s)\*(?!\*)", r"<em>\1</em>", s)  # italic
        s = re.sub(r"`([^`]+)`", r"<code>\1</code>", s)           # inline code
        s = re.sub(r"\[([^\]]+)\]\((https?://[^\s)]+)\)", r'<a href="\2" target="_blank">\1</a>', s)  # links
        return s

    lines = md.split("\n")
    out, in_ul, in_ol, in_table = [], False, False, False

    def close_lists():
        nonlocal in_ul, in_ol
        if in_ul:
            out.append("</ul>")
            in_ul = False
        if in_ol:
            out.append("</ol>")
            in_ol = False

    for ln in lines:
        raw = ln.strip()

        # Horizontal rule
        if re.match(r"^---+$", raw):
            close_lists(); out.append("<hr/>"); continue

        # Headings
        m = re.match(r"^(#{1,6})\s+(.+)$", raw)
        if m:
            close_lists()
            level = len(m.group(1))
            out.append(f"<h{level}>{fmt_inline(m.group(2))}</h{level}>")
            continue

        # Table rows
        if raw.startswith("|") and raw.endswith("|"):
            cells = [fmt_inline(c.strip()) for c in raw.strip("|").split("|")]
            if not in_table:
                out.append("<table><tbody>")
                in_table = True
            out.append("<tr>" + "".join(f"<td>{c}</td>" for c in cells) + "</tr>")
            continue
        else:
            if in_table:
                out.append("</tbody></table>")
                in_table = False

        # Ordered list
        m = re.match(r"^\d+[.)]\s+(.+)$", raw)
        if m:
            if in_ul: out.append("</ul>"); in_ul = False
            if not in_ol: out.append("<ol>"); in_ol = True
            out.append(f"<li>{fmt_inline(m.group(1))}</li>")
            continue

        # Unordered list
        m = re.match(r"^[-*•–]\s+(.+)$", raw)
        if m:
            if in_ol: out.append("</ol>"); in_ol = False
            if not in_ul: out.append("<ul>"); in_ul = True
            out.append(f"<li>{fmt_inline(m.group(1))}</li>")
            continue

        # Paragraphs
        if raw:
            close_lists()
            out.append(f"<p>{fmt_inline(raw)}</p>")
        else:
            close_lists()

    close_lists()
    if in_table:
        out.append("</tbody></table>")

    return "\n".join(out)

def render_chat_history(messages: List[Dict[str, str]]) -> None:
    """Assistant left; User right (with emoji icons)."""
    assistant_icon = "🤖"
    user_icon = "🕵️‍♂️"
    for m in messages:
        role = (m.get("role", "assistant") or "assistant").lower()
        content = (m.get("content", "") or "").strip()
        if not content:
            continue
        if role == "user":
            inner_html = html.escape(content).replace("\n", "<br>")
            st.markdown(
                f"""
                <div class="msg-row msg-right">
                  <div class="msg-bubble msg-user">{inner_html}</div>
                  <div class="msg-icon">{user_icon}</div>
                </div>
                """,
                unsafe_allow_html=True,
            )
        else:
            inner_html = format_llm_reply_to_html(content)
            st.markdown(
                f"""
                <div class="msg-row msg-left">
                  <div class="msg-icon">{assistant_icon}</div>
                  <div class="msg-bubble msg-assistant">{inner_html}</div>
                </div>
                """,
                unsafe_allow_html=True,
            )

# =======================
# ❖ Sidebar (logo → nav → country search → tools)
# =======================
def render_sidebar_home() -> None:
    with st.sidebar:
        # 1) Logo
        show_logo(center=True, max_width="180px")

        # 2) Navigation
        st.markdown("### Navigation")
        try:
            st.page_link("Home.py", label="🏠 Chat", icon=None)
            st.page_link("pages/Support.py", label="❓ Support", icon=None)
        except Exception:
            st.write("• Chat")
        st.markdown("---")

        # 3) Country search
        st.markdown("#### Country search")
        country_name = st.selectbox(
            "Search or select a country",
            options=[n for n, _ in COUNTRIES],
            index=None,
            placeholder="Type to search…",
        )
        if country_name:
            selected = find_country_by_name(country_name)
            if selected:
                n, iso = selected
                st.session_state["selected_country_name"] = n
                st.session_state["selected_country_iso"] = iso
                st.query_params.clear()
                st.query_params["country"] = iso
                st.session_state["__go_country__"] = True
                st.rerun()

        st.markdown("---")

        # 4) Tools — Clear conversation
        if st.button("Clear conversation", use_container_width=True):
            st.session_state[SK_MSGS] = [
                {"role": "assistant", "content": "Hi! I am TomTom's Tax Assistant. How can I help today?"}
            ]
            st.rerun()

# =======================
# ❖ Top-level redirect guard
# =======================
if st.session_state.get("__go_country__") or st.query_params.get("country"):
    if "__go_country__" in st.session_state:
        del st.session_state["__go_country__"]
    try:
        st.switch_page("pages/_Country.py")
    except Exception:
        st.rerun()

# =======================
# ❖ Chat UI
# =======================
def chat_ui() -> None:
    render_sidebar_home()

    st.title(APP_TITLE)
    st.markdown("---")
    st.header("📖 Chat Assistant")

    st.session_state.setdefault(SK_MSGS, [{"role": "assistant", "content": "Hi! I am TomTom's Tax Assistant. How can I help today?"}])
    st.session_state.setdefault(ATTACH_CTX_KEY, [])

    # Chat area
    st.markdown('<div class="chat-margin-container">', unsafe_allow_html=True)
    render_chat_history(st.session_state[SK_MSGS])
    st.markdown('</div>', unsafe_allow_html=True)  # close before bottom widgets

    # Bottom-left floating "Attach" popover
    st.markdown('<div class="attach-anchor">', unsafe_allow_html=True)
    with st.popover("➕ Attach", use_container_width=False):
        st.caption("Upload files to enrich context (text is extracted and appended).")
        up_files = st.file_uploader(
            "Select file(s)",
            type=[
                "pdf","docx","txt","md","csv","tsv","json","yaml","yml","log",
                "xlsx","pptx","rtf","html","htm","xml",
                "png","jpg","jpeg","bmp","gif","webp","tif","tiff"
            ],
            accept_multiple_files=True,
            label_visibility="collapsed",
            key="tt_attach_uploader",
        )
        c1, c2 = st.columns(2)
        with c1:
            if st.button("Add to context", use_container_width=True, key="tt_add_ctx_btn"):
                added_any = False
                if up_files:
                    for f in up_files:
                        try:
                            txt = extract_text_from_upload(f)
                            if txt:
                                st.session_state[ATTACH_CTX_KEY].append(txt)
                                added_any = True
                        except Exception:
                            pass
                if added_any:
                    joined = "\n\n".join(st.session_state[ATTACH_CTX_KEY])
                    if len(joined) > MAX_TOTAL_ATTACHMENT_CHARS:
                        joined = joined[:MAX_TOTAL_ATTACHMENT_CHARS] + "\n...[attachments truncated]"
                        st.session_state[ATTACH_CTX_KEY] = [joined]
                    st.success("Attachment(s) added.")
                    st.rerun()
                else:
                    st.warning("No readable content extracted.")
        with c2:
            if st.button("Clear attachments", use_container_width=True, key="tt_clear_ctx_btn"):
                st.session_state.pop(ATTACH_CTX_KEY, None)
                st.rerun()
    st.markdown('</div>', unsafe_allow_html=True)

    # Chat input LAST, so it stays at the bottom
    prompt = st.chat_input("Type your message and press enter", key="chat_prompt")
    if prompt and prompt.strip():
        st.session_state[SK_MSGS].append({"role": "user", "content": prompt.strip()})
        with st.spinner("Thinking…"):
            reply_raw = get_llm_response(prompt.strip())
        st.session_state[SK_MSGS].append({"role": "assistant", "content": reply_raw})

        # keep context window short
        if len(st.session_state[SK_MSGS]) > MAX_CONTEXT_MESSAGES:
            st.session_state[SK_MSGS] = st.session_state[SK_MSGS][-MAX_CONTEXT_MESSAGES:]
        st.rerun()

# =======================
# ❖ App entry
# =======================
def main() -> None:
    # Preload KB once per session
    if KB_STATE_KEY not in st.session_state:
        st.session_state[KB_STATE_KEY] = build_bootstrap_kb(BOOTSTRAP_DOCS_DIR)
    chat_ui()

if __name__ == "__main__":
    main()
